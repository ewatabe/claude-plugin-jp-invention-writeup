#!/usr/bin/env python3
"""発明説明資料PPTXビルダ。

`work/invention-explainer-draft.md` を読み、サンプル `~/patent/sample/発明説明資料.pptx`
準拠の構成（表紙＋アジェンダ＋5セクション）でPPTXを生成する。

入力markdownのルール:
- `# 表紙`: 表紙スライド（- key: value 形式の箇条書きから日付・所属・発明者・案件名を抽出）
- `# アジェンダ`: アジェンダスライド
- `# 1. ...` 〜 `# 5. ...`: セクション扉スライド
- `## ...`: 通常スライド（タイトルは ## の見出し）
- `- ...`: スライドの bullet
- `![alt](path)`: スライドに画像を挿入

Usage:
    python3 build_explainer_pptx.py \\
        --input ~/patent/<案件名>/work/invention-explainer-draft.md \\
        --output ~/patent/<案件名>/output/発明説明資料.pptx
"""
from __future__ import annotations

import argparse
import datetime
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
JP_FONT = "Yu Gothic"


def in_(v):
    return Inches(v)


def parse_markdown(md_text: str):
    """markdownを (slide_type, title, content_items) のリストに変換。"""
    lines = md_text.splitlines()
    slides = []
    current = None
    in_h1 = False

    for line in lines:
        h1 = re.match(r"^# (.+)$", line)
        h2 = re.match(r"^## (.+)$", line)
        bullet = re.match(r"^- (.+)$", line)
        image = re.match(r"^!\[[^\]]*\]\(([^)]+)\)$", line)

        if h1:
            if current:
                slides.append(current)
            title = h1.group(1).strip()
            stype = (
                "cover" if "表紙" in title
                else "agenda" if "アジェンダ" in title
                else "section"
            )
            current = {"type": stype, "title": title, "items": [], "images": []}
            in_h1 = True
        elif h2:
            if current:
                slides.append(current)
            current = {"type": "content", "title": h2.group(1).strip(), "items": [], "images": []}
            in_h1 = False
        elif bullet:
            if current is not None:
                current["items"].append(bullet.group(1).strip())
        elif image:
            if current is not None:
                current["images"].append(image.group(1).strip())

    if current:
        slides.append(current)
    return slides


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, align=PP_ALIGN.LEFT, color=(0, 0, 0)):
    tb = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = JP_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    return tb


def render_cover(prs, slide_data):
    """表紙スライド。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 日付・所属・発明者・案件名を items から拾う
    meta = {}
    for item in slide_data["items"]:
        m = re.match(r"^([^:]+):\s*(.+)$", item)
        if m:
            meta[m.group(1).strip()] = m.group(2).strip()

    today = datetime.date.today()
    date_str = meta.get("日付", "").replace("{{YYYY/M/D}}", "")
    if not date_str or "YYYY" in date_str:
        date_str = f"{today.year}/{today.month}/{today.day}"

    add_textbox(slide, 1.0, 1.0, 11, 0.6, date_str, size=14)
    add_textbox(slide, 1.0, 1.8, 11, 0.6, meta.get("所属", ""), size=14)
    add_textbox(slide, 1.0, 3.0, 11, 1.0, meta.get("案件名", "") or "発明説明資料資料", size=32, bold=True)
    add_textbox(slide, 1.0, 5.5, 11, 0.6, "発明者: " + meta.get("発明者", ""), size=14)


def render_agenda(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 1.0, 0.5, 11, 0.8, "アジェンダ", size=28, bold=True)
    y = 1.8
    for item in slide_data["items"]:
        add_textbox(slide, 1.5, y, 10, 0.5, "・ " + item, size=18)
        y += 0.55


def render_section(prs, slide_data):
    """セクション扉スライド（背景色＋大タイトル）。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 1.0, 2.8, 11, 1.5, slide_data["title"], size=40, bold=True)


def render_content(prs, slide_data):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.6, 0.4, 12, 0.7, slide_data["title"], size=24, bold=True)

    y = 1.5
    for item in slide_data["items"]:
        add_textbox(slide, 1.0, y, 11.5, 0.5, "・ " + item, size=16)
        y += 0.55

    for img_path in slide_data["images"]:
        p = Path(img_path).expanduser()
        if p.exists():
            slide.shapes.add_picture(str(p), in_(1.0), in_(y + 0.2), height=in_(SLIDE_H_IN - y - 0.7))
            break  # 1スライド1画像


def build(input_md: Path, output_path: Path):
    md_text = input_md.read_text(encoding="utf-8")
    slides = parse_markdown(md_text)

    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))

    renderers = {
        "cover": render_cover,
        "agenda": render_agenda,
        "section": render_section,
        "content": render_content,
    }

    for s in slides:
        renderers.get(s["type"], render_content)(prs, s)
        print(f"  added: [{s['type']}] {s['title']}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"OK: {output_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()
    build(Path(args.input).expanduser(), Path(args.output).expanduser())


if __name__ == "__main__":
    main()
