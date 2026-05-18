#!/usr/bin/env python3
"""図面PPTXビルダ。

`work/figure-specs/figure-*.json` を順に読み、サンプル `~/patent/sample/図面1.pptx`
準拠のヘッダ（【書類名】図面 / 【図N】 / 図 N）付きスライドを1枚ずつ追加して
`output/図面.pptx` を生成する。

Usage:
    python3 build_figure_pptx.py \\
        --specs ~/patent/<案件名>/work/figure-specs/ \\
        --output ~/patent/<案件名>/output/図面.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor

# pptx_helpers から set_run_font_full をインポート
sys.path.insert(0, str(Path(__file__).parent))
from pptx_helpers import set_run_font_full  # noqa: E402

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
JP_FONT = "Yu Gothic UI"


def in_(v):
    return Inches(v)


def add_header(slide, figure_number):
    """JPO実務に合わせたヘッダ（【書類名】図面、【図N】、図 N）を配置。"""
    # 【書類名】 図面 (左上)
    tb1 = slide.shapes.add_textbox(in_(0.3), in_(0.2), in_(2.5), in_(0.3))
    p = tb1.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "【書類名】 図面"
    set_run_font_full(r, JP_FONT)
    r.font.size = Pt(11)

    # 【図N】 (左上、書類名の下)
    tb2 = slide.shapes.add_textbox(in_(0.3), in_(0.55), in_(2.0), in_(0.3))
    p = tb2.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"【図{figure_number}】"
    set_run_font_full(r, JP_FONT)
    r.font.size = Pt(11)

    # 図 N (中央上、大きめ)
    tb3 = slide.shapes.add_textbox(in_(SLIDE_W_IN / 2 - 0.5), in_(0.2), in_(1.0), in_(0.5))
    p = tb3.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"図 {figure_number}"
    set_run_font_full(r, JP_FONT)
    r.font.size = Pt(18)
    r.font.bold = True


def add_rectangle(slide, comp):
    """構成要素の四角形と内部ラベル、近傍の参照符号を配置。"""
    x, y = comp["pos"]
    w, h = comp["size"]
    shape_type = comp.get("shape", "rectangle")

    msoshape = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "cylinder": MSO_SHAPE.CAN,
        "diamond": MSO_SHAPE.DIAMOND,
        "rhombus": MSO_SHAPE.DIAMOND,
    }.get(shape_type, MSO_SHAPE.RECTANGLE)

    shp = slide.shapes.add_shape(msoshape, in_(x), in_(y), in_(w), in_(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shp.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    shp.line.width = Pt(1)

    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = comp.get("label", "")
    set_run_font_full(r, JP_FONT)
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # 参照符号（別テキストボックスで配置）
    ref = comp.get("ref", "")
    if ref:
        anchor = comp.get("ref_anchor", "top_left")
        if anchor == "top_left":
            rx, ry = x - 0.3, y - 0.25
        elif anchor == "top_right":
            rx, ry = x + w + 0.05, y - 0.25
        elif anchor == "bottom_left":
            rx, ry = x - 0.3, y + h
        elif anchor == "bottom_right":
            rx, ry = x + w + 0.05, y + h
        else:
            rx, ry = x - 0.3, y - 0.25

        rtb = slide.shapes.add_textbox(in_(max(rx, 0)), in_(max(ry, 0)), in_(0.4), in_(0.25))
        rp = rtb.text_frame.paragraphs[0]
        rr = rp.add_run()
        rr.text = ref
        set_run_font_full(rr, JP_FONT)
        rr.font.size = Pt(11)
        rr.font.bold = True

    return shp


def add_arrow(slide, frm_comp, to_comp, direction):
    """from → to の中心同士を結ぶ矢印を1本追加。

    direction: right/left/up/down/bidirectional
    """
    fx, fy = frm_comp["pos"]
    fw, fh = frm_comp["size"]
    tx, ty = to_comp["pos"]
    tw, th = to_comp["size"]

    # 端点を求める（最も近い辺の中点を結ぶ簡易ロジック）
    fcx, fcy = fx + fw / 2, fy + fh / 2
    tcx, tcy = tx + tw / 2, ty + th / 2

    if abs(tcx - fcx) >= abs(tcy - fcy):
        # 水平方向に支配的
        if tcx > fcx:
            sx, sy = fx + fw, fcy
            ex, ey = tx, tcy
        else:
            sx, sy = fx, fcy
            ex, ey = tx + tw, tcy
    else:
        if tcy > fcy:
            sx, sy = fcx, fy + fh
            ex, ey = tcx, ty
        else:
            sx, sy = fcx, fy
            ex, ey = tcx, ty + th

    # 矢印図形（適切なMSO_SHAPE）
    msoshape = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "bidirectional": MSO_SHAPE.LEFT_RIGHT_ARROW,
    }.get(direction, MSO_SHAPE.RIGHT_ARROW)

    # 簡略のため、from と to の中点を結ぶ "block arrow" を1本描く
    x = min(sx, ex)
    y = min(sy, ey)
    w = max(abs(ex - sx), 0.3)
    h = max(abs(ey - sy), 0.25)
    if direction in ("up", "down"):
        w = 0.3
        x = (sx + ex) / 2 - 0.15
    else:
        h = 0.25
        y = (sy + ey) / 2 - 0.125

    shp = slide.shapes.add_shape(msoshape, in_(x), in_(y), in_(w), in_(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shp.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    return shp


def build_block_diagram(prs, spec):
    """構成図 / ブロック図 を1スライド追加。"""
    blank = prs.slide_layouts[-1]  # 完全空白を末尾レイアウトから取る
    # "白紙" を含む層があればそれを優先
    for layout in prs.slide_layouts:
        if "白紙" in layout.name or "Blank" in layout.name or layout.name == "":
            blank = layout
            break
    slide = prs.slides.add_slide(blank)
    add_header(slide, spec["figure_number"])

    comps_by_id = {}
    for comp in spec.get("components", []):
        add_rectangle(slide, comp)
        comps_by_id[comp["id"]] = comp

    for arrow in spec.get("arrows", []):
        frm = comps_by_id.get(arrow["from"])
        to = comps_by_id.get(arrow["to"])
        if frm and to:
            add_arrow(slide, frm, to, arrow.get("direction", "right"))


def build_table(prs, spec):
    """表形式の図1スライド。spec.headers と spec.rows を表として配置。"""
    blank = prs.slide_layouts[-1]
    for layout in prs.slide_layouts:
        if "白紙" in layout.name or "Blank" in layout.name or layout.name == "":
            blank = layout
            break
    slide = prs.slides.add_slide(blank)
    add_header(slide, spec["figure_number"])

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers or not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 ヘッダ行

    # 配置: 表幅は 11 inch、左右マージン 1.15 inch
    table_w = SLIDE_W_IN - 2.3
    table_h = min(0.5 + n_rows * 0.45, SLIDE_H_IN - 2.5)
    left = (SLIDE_W_IN - table_w) / 2
    top = 1.3

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, in_(left), in_(top), in_(table_w), in_(table_h)
    )
    tbl = table_shape.table

    # ヘッダ
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                set_run_font_full(run, JP_FONT)
                run.font.size = Pt(11)
                run.font.bold = True

    # データ行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    set_run_font_full(run, JP_FONT)
                    run.font.size = Pt(10)


def build_image(prs, spec):
    """画像埋め込みスライド。spec.image_path で指定された画像を中央配置。"""
    blank = prs.slide_layouts[-1]
    for layout in prs.slide_layouts:
        if "白紙" in layout.name or "Blank" in layout.name or layout.name == "":
            blank = layout
            break
    slide = prs.slides.add_slide(blank)
    add_header(slide, spec["figure_number"])

    img_path = Path(spec["image_path"]).expanduser()
    if not img_path.exists():
        print(f"  WARN: image not found: {img_path}", file=sys.stderr)
        return

    # 中央配置、ヘッダの下に最大サイズで
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size

    max_h = SLIDE_H_IN - 1.5  # ヘッダ分を引く
    max_w = SLIDE_W_IN - 1.0
    h = max_h
    w = h * iw / ih
    if w > max_w:
        w = max_w
        h = w * ih / iw

    x = (SLIDE_W_IN - w) / 2
    y = 1.0 + (max_h - h) / 2
    slide.shapes.add_picture(str(img_path), in_(x), in_(y), width=in_(w), height=in_(h))


def build_flowchart(prs, spec):
    """フローチャート1スライド。steps を縦に並べ、矢印で接続。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header(slide, spec["figure_number"])

    steps = spec.get("steps", [])
    if not steps:
        return

    # 縦に並べる：上から下、各ステップ x=中央
    cx = SLIDE_W_IN / 2
    box_w, box_h = 3.0, 0.7
    gap = 0.5
    start_y = 1.2
    positions = {}
    for i, step in enumerate(steps):
        y = start_y + i * (box_h + gap)
        shape_id = step["id"]
        shape_type = step.get("type", "process")
        shape_kind = {
            "start": "rounded_rectangle",
            "end": "rounded_rectangle",
            "process": "rectangle",
            "decision": "diamond",
        }.get(shape_type, "rectangle")

        comp = {
            "id": shape_id,
            "label": step["label"],
            "ref": step.get("ref", ""),
            "shape": shape_kind,
            "pos": [cx - box_w / 2, y],
            "size": [box_w, box_h],
            "ref_anchor": "top_left",
        }
        add_rectangle(slide, comp)
        positions[shape_id] = comp

    for tr in spec.get("transitions", []):
        frm = positions.get(tr["from"])
        to = positions.get(tr["to"])
        if frm and to:
            add_arrow(slide, frm, to, "down")


def _add_full_slide_image(prs, img_path):
    """1スライドに 1枚の画像をフル配置（ヘッダ/タイトル等もすべて画像側で完結）。"""
    blank = prs.slide_layouts[-1]
    for layout in prs.slide_layouts:
        if "白紙" in layout.name or "Blank" in layout.name or layout.name == "":
            blank = layout
            break
    slide = prs.slides.add_slide(blank)
    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    # スライドサイズに最大フィット（アスペクト比保持）
    sw = SLIDE_W_IN
    sh = SLIDE_H_IN
    scale = min(sw / (iw / 96), sh / (ih / 96))
    # 96 dpi は仮の換算。実際は画像のアスペクト比のみ重要
    iw_aspect = iw / ih
    if iw_aspect >= sw / sh:
        w = sw
        h = sw / iw_aspect
    else:
        h = sh
        w = sh * iw_aspect
    x = (sw - w) / 2
    y = (sh - h) / 2
    slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def build_pptx(specs_dir: Path, output_path: Path, rendered_dir: Path = None):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))

    json_files = sorted(specs_dir.glob("figure-*.json"))
    if not json_files:
        print(f"WARN: figure-*.json が {specs_dir} に見つかりません", file=sys.stderr)
        sys.exit(1)

    # rendered モード: rendered_dir が指定されていれば JSON の figure_type を無視して
    # 既にレンダ済みの PNG (figure-NN.png) を全スライドのフルキャンバスとして埋め込む
    if rendered_dir is not None:
        rendered_dir = Path(rendered_dir).expanduser()
        for f in json_files:
            spec = json.loads(f.read_text(encoding="utf-8"))
            n = spec.get("figure_number")
            png = rendered_dir / f"figure-{n:02d}.png"
            if not png.exists():
                print(f"  WARN: rendered PNG not found: {png}", file=sys.stderr)
                continue
            _add_full_slide_image(prs, png)
            print(f"  added figure {n}: {spec.get('title', '')} (rendered)")
    else:
        for f in json_files:
            spec = json.loads(f.read_text(encoding="utf-8"))
            ftype = spec.get("figure_type", "block_diagram")
            if ftype == "flowchart":
                build_flowchart(prs, spec)
            elif ftype == "table":
                build_table(prs, spec)
            elif ftype == "image":
                build_image(prs, spec)
            else:
                build_block_diagram(prs, spec)
            print(f"  added figure {spec.get('figure_number')}: {spec.get('title', '')}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"OK: {output_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--specs", required=True, help="figure-*.json があるディレクトリ")
    ap.add_argument("--output", required=True, help="出力pptxパス")
    ap.add_argument("--rendered-dir",
                    help="HTMLレンダ済みPNG（figure-NN.png）のディレクトリ。指定するとフル画像埋込モードに切替")
    args = ap.parse_args()
    rd = Path(args.rendered_dir).expanduser() if args.rendered_dir else None
    build_pptx(Path(args.specs).expanduser(), Path(args.output).expanduser(), rendered_dir=rd)


if __name__ == "__main__":
    main()
