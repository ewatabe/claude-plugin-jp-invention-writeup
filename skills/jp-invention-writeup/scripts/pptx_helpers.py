"""python-pptx の共通ヘルパー関数群。

発明説明資料PPTXの組立で繰り返し使う処理を集約：
- レイアウト選択（'タイトルのみ' / cover / blank の正しい判別）
- テキスト置換（先頭runのフォント・色・サイズを保持）
- 噴き出し（ROUNDED_RECTANGLE）+ leader line（connector）
- 画像 + 複数噴き出しの一括配置
- スライド並べ替え（XML操作）
- フォント統一（Yu Gothic UI 等を全テキストランに一括適用）

使い方:
    from pptx_helpers import (
        add_content_slide, replace_text_preserving_style,
        add_speech_bubble, add_image_with_bubbles,
        reorder_slides, unify_font,
    )
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

DEFAULT_FONT = "Yu Gothic UI"


# ============================================================
# 重要：OOXML フォント指定（latin / ea / cs）
# ============================================================
# python-pptx の `run.font.name = ...` は latin typeface のみ設定する。
# 日本語（East Asian）文字は ea typeface のフォールバックを使うため、
# テーマ既定フォント（游ゴシック等）が描画されてしまう。
# PowerPoint上でフォント名は "Yu Gothic UI" と表示されるが、
# 実描画は別フォント、というズレが起きる。
# set_run_font_full() で latin/ea/cs を全て設定し、ズレを解消する。

def set_run_font_full(run, font_name):
    """テキストランに対し、latin / ea / cs の全 typeface を設定。
    日本語文字も Yu Gothic UI で表示されることを保証する。
    """
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        # 既存タグを削除
        for el in rPr.findall(qn(f"a:{tag}")):
            rPr.remove(el)
        # 新規追加
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", font_name)

# 噴き出しの既定スタイル（淡い黄色＋金茶ボーダー、テキストは黒）
DEFAULT_BUBBLE_FILL = (0xFF, 0xF5, 0xCE)
DEFAULT_BUBBLE_BORDER = (0xCC, 0xA8, 0x00)
DEFAULT_BUBBLE_TEXT = (0x00, 0x00, 0x00)

# サンプルユーザー名（実施例で使用）
SAMPLE_USER_NAME = "山田 花子"
SAMPLE_USER_NAME_KANJI = "山田花子"


# ============================================================
# レイアウト選択
# ============================================================
def find_layout(prs, *needles, fallback_index=-1):
    """layout.name が needle と一致するレイアウトを返す。
    優先度：完全一致 > 前方一致（needle で始まる） > 部分一致。
    複数 needle を渡した場合は最初にマッチした needle のレイアウトを返す。
    重要：'タイトルのみ' を探したとき '4_タイトルのみ'（表紙）を誤って返さないために完全一致を優先する。
    """
    for needle in needles:
        # 完全一致
        for layout in prs.slide_layouts:
            if layout.name == needle:
                return layout
        # 前方一致
        for layout in prs.slide_layouts:
            if layout.name.startswith(needle):
                return layout
        # 部分一致
        for layout in prs.slide_layouts:
            if needle in layout.name:
                return layout
    return prs.slide_layouts[fallback_index]


def add_content_slide(prs, title_text=None, layout_name_contains=("タイトルのみ",), font_name=DEFAULT_FONT):
    """通常コンテンツ用レイアウトでスライドを追加。タイトルをセット。

    重要：`4_白紙` (closing logo) は使わない。'タイトルのみ' を使う。
    """
    layout = find_layout(prs, *layout_name_contains, fallback_index=2)
    slide = prs.slides.add_slide(layout)
    if title_text:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.placeholder_format is not None:
                if shape.placeholder_format.type == 1:  # TITLE
                    shape.text_frame.text = title_text
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            set_run_font_full(r, font_name)
                    break
    return slide


# ============================================================
# テキスト置換（スタイル保持）
# ============================================================
def find_shape_by_substring(slide, needle):
    """テキストフレーム内に needle を含む最初の shape を返す。"""
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    return None


def replace_text_preserving_style(shape, new_text, *, font_name=None):
    """テキストフレームの中身を新テキストに置換。
    先頭runのフォント名・サイズ・bold・色を全段落に引き継ぐ。
    font_name を指定するとさらに上書き。

    重要：title placeholder 等のテキストを更新する際は **必ずこの関数経由で**
    置換すること。`for run in para.runs: run.text = "..."` のように
    各runを個別に書き換えると、元テキストが複数runに分かれている場合に
    同じ新テキストが複数runに増殖して **重複表示** されるバグになる。
    """
    tf = shape.text_frame
    fn, fs, fb, fc = None, None, None, None
    if tf.paragraphs and tf.paragraphs[0].runs:
        r0 = tf.paragraphs[0].runs[0]
        fn = r0.font.name
        fs = r0.font.size
        fb = r0.font.bold
        try:
            fc = r0.font.color.rgb
        except Exception:
            fc = None
    tf.clear()
    for i, line in enumerate(str(new_text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_run_font_full(run, font_name or fn or DEFAULT_FONT)
        if fs:
            run.font.size = fs
        if fb is not None:
            run.font.bold = fb
        if fc:
            try:
                run.font.color.rgb = fc
            except Exception:
                pass


# ============================================================
# テキストボックス / 噴き出し / 線
# ============================================================
def add_textbox(slide, x, y, w, h, text, *,
                size=14, bold=False, align=PP_ALIGN.LEFT,
                color=(0, 0, 0), font_name=DEFAULT_FONT, wrap=True):
    """シンプルなテキストボックスを追加。x,y,w,h は inch 単位。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        set_run_font_full(r, font_name)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)
    return tb


def add_speech_bubble(slide, x, y, w, h, text, *,
                      fill=DEFAULT_BUBBLE_FILL, border=DEFAULT_BUBBLE_BORDER,
                      text_color=DEFAULT_BUBBLE_TEXT,
                      font_size=10, font_name=DEFAULT_FONT,
                      anchor=None):
    """ROUNDED_RECTANGULAR_CALLOUT（角丸吹き出し）を追加。
    anchor=(x, y) が指定された場合、吹き出しのしっぽが
    そのスライド座標を指すよう adj 値を調整する。
    anchor=None なら ROUNDED_RECTANGLE（しっぽなし）にフォールバック。
    """
    if anchor is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                                     Inches(x), Inches(y), Inches(w), Inches(h))
        # しっぽの先端を anchor へ。adj1=tail x (fraction of width from shape left)、adj2=tail y (fraction of height from top)
        try:
            shp.adjustments[0] = (anchor[0] - x) / w
            shp.adjustments[1] = (anchor[1] - y) / h
        except (IndexError, AttributeError):
            pass
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*border)
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        set_run_font_full(r, font_name)
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(*text_color)
    return shp


def add_leader_line(slide, x1, y1, x2, y2, *,
                    color=DEFAULT_BUBBLE_BORDER, width=1.25):
    """噴き出しからターゲットへの引出線（直線）。"""
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    return line


def add_image_with_bubbles(slide, img_path, img_top, img_max_h, bubbles, *,
                           slide_w_in=13.333, margin=0.5,
                           bubble_kwargs=None, leader_kwargs=None):
    """画像をスライド中央に配置し、その周囲に複数の吹き出し（角丸callout）を一括配置。

    吹き出しは ROUNDED_RECTANGULAR_CALLOUT（角丸＋しっぽ付き）で、
    しっぽが anchor 座標を指すよう自動調整される。

    Args:
        img_path: 画像パス
        img_top: 画像の上端 (inch)
        img_max_h: 画像の最大高さ (inch)
        bubbles: [(x, y, w, h, text, anchor)] のリスト。
                 anchor は (target_x, target_y) または None
        slide_w_in: スライドの横幅 (inch)
        margin: 画像周りの左右マージン (inch)
        leader_kwargs: 後方互換用、無視される
    Returns:
        (img_x, img_y, img_w, img_h): 配置された画像の座標とサイズ
    """
    from PIL import Image as PILImage

    img_path = Path(str(img_path)).expanduser()
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    h = img_max_h
    w = h * iw / ih
    max_w = slide_w_in - 2 * margin
    if w > max_w:
        w = max_w
        h = w * ih / iw
    img_x = (slide_w_in - w) / 2
    slide.shapes.add_picture(str(img_path), Inches(img_x), Inches(img_top),
                              width=Inches(w), height=Inches(h))

    bk = bubble_kwargs or {}
    for entry in bubbles:
        bx, by, bw, bh, text, anchor = entry
        # callout shape のしっぽが anchor を指す
        add_speech_bubble(slide, bx, by, bw, bh, text, anchor=anchor, **bk)
    return img_x, img_top, w, h


# ============================================================
# スライド並べ替え
# ============================================================
def reorder_slides(prs, desired_order):
    """スライドを desired_order の0-indexedインデックス順に並べ替える。

    Args:
        prs: Presentation オブジェクト
        desired_order: 並べ替え後のスライド順を表す元index のリスト
                       省略されたindexはpptxから除外される
    """
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    for sld_id in sld_ids:
        xml_slides.remove(sld_id)
    for idx in desired_order:
        xml_slides.append(sld_ids[idx])


# ============================================================
# フォント統一
# ============================================================
def unify_font(prs, font_name=DEFAULT_FONT):
    """全スライドの全テキストランに同じフォント名を適用する。
    latin / ea / cs の全 typeface を設定するため日本語も正しく表示される。
    """
    n_runs = 0
    for slide in prs.slides:
        n_runs += _apply_font_to_shapes(slide.shapes, font_name)
    return n_runs


def _apply_font_to_shapes(shapes, font_name):
    n_runs = 0
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    set_run_font_full(run, font_name)
                    n_runs += 1
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            set_run_font_full(run, font_name)
                            n_runs += 1
        # グループ図形は再帰
        try:
            n_runs += _apply_font_to_shapes(shape.shapes, font_name)
        except AttributeError:
            pass
    return n_runs


# ============================================================
# ページ番号管理（テンプレ内のハードコード番号を更新）
# ============================================================
def update_hardcoded_page_numbers(prs, *, start_page=1, only_short_text=True):
    """各スライドに含まれる「1-3桁の数字だけのテキストボックス」を
    現在のスライド位置のページ番号で上書きする。

    元テンプレで個別テキストボックスに番号がハードコードされている場合の救済。

    Args:
        start_page: スライド1に振る番号（既定 1）
        only_short_text: True なら text が "数字のみ" のものだけ対象。
                         False は危険（他の用途の番号も上書きしうる）
    """
    import re
    updated = 0
    for i, slide in enumerate(prs.slides):
        page = start_page + i
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if only_short_text and not re.fullmatch(r"\d{1,3}", text):
                continue
            # フォント情報を保持しつつ上書き
            replace_text_preserving_style(shape, str(page))
            updated += 1
            break  # 1スライドに1番号だけ対象
    return updated
