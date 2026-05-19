#!/usr/bin/env python3
"""Rich slide deck builder — 役員説明レベルの発明説明資料 PPTX ビルダ.

HTML+CSS で個別スライドを記述 → Chrome ヘッドレスで 1920x1080 PNG レンダ
→ python-pptx で 1 スライド 1 画像のフルブリード配置で PPTX 化する。

ファイル順序: figures-rendered/ 配下の *.png をファイル名昇順で並べる。
ファイル名規約: NN-name.png（例: 01-cover.png, 04b-extra.png）

使い方:
    python3 build_rich_pptx.py \\
        --rendered-dir ~/patent/<案件名>/work/figures-rendered/ \\
        --output ~/patent/<案件名>/output/発明説明資料.pptx

事前にレンダ済 PNG を用意しておく必要がある。レンダは render_all_slides.sh
ヘルパ（同ディレクトリ）または Chrome ヘッドレスを直接呼ぶ。
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# 16:9 1920x1080 を 13.333 x 7.5 inch スライドへフルブリード配置
W_IN = 13.333
H_IN = 7.5


def build(rendered_dir: Path, output_path: Path) -> None:
    pngs = sorted(rendered_dir.glob("*.png"))
    if not pngs:
        print(f"ERROR: no PNG files in {rendered_dir}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Emu(int(W_IN * 914400))
    prs.slide_height = Emu(int(H_IN * 914400))

    blank = prs.slide_layouts[6]  # 完全白紙レイアウト
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png),
            left=Emu(0),
            top=Emu(0),
            width=Emu(int(W_IN * 914400)),
            height=Emu(int(H_IN * 914400)),
        )
        print(f"  added: {png.name}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"OK: {output_path} ({len(pngs)} slides)")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--rendered-dir",
        required=True,
        help="HTML レンダ済 PNG (NN-name.png) のディレクトリ",
    )
    ap.add_argument("--output", required=True, help="出力 PPTX のパス")
    args = ap.parse_args()

    build(Path(args.rendered_dir).expanduser(), Path(args.output).expanduser())


if __name__ == "__main__":
    main()
